"""Generate TyvaTransit diploma presentation from college template."""
from __future__ import annotations

import copy
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = Path(r"C:\Users\Bailak\Downloads\Telegram Desktop\Структура презентации.pptx")
OUTPUT = Path(r"D:\VScode_Project\TyvaTransit\TyvaTransit_презентация.pptx")
SCREENSHOTS = Path(r"D:\VScode_Project\TyvaTransit\presentation_screenshots")

# Placeholders — замените на свои данные перед защитой
STUDENT = "ФИО студента"
GROUP = "группа ___"
SUPERVISOR = "руководитель: ФИО"
DEMO_URL = "https://ваш-домен.up.railway.app"


def set_title(slide, text: str) -> None:
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 1:  # TITLE
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(32)
            p.font.bold = True
            return
    for shape in slide.shapes:
        if shape.has_text_frame and "Заголовок" in shape.name:
            shape.text_frame.text = text
            return


def add_body_text(slide, lines: list[str], *, left=None, top=None, width=None, height=None) -> None:
    left = left or Inches(0.92)
    top = top or Inches(2.0)
    width = width or Inches(11.5)
    height = height or Inches(4.8)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.space_after = Pt(6)
        if line.endswith(":") or (line and not line.startswith("•") and i == 0 and len(lines) > 1):
            p.font.bold = True


def add_bullets(slide, items: list[str], **kwargs) -> None:
    lines = []
    for item in items:
        lines.append(f"• {item}")
    add_body_text(slide, lines, **kwargs)


def add_image(slide, image_path: Path, *, left, top, width) -> None:
    slide.shapes.add_picture(str(image_path), left, top, width=width)


def fill_title_slide(slide) -> None:
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 3:  # CENTER_TITLE
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = (
                "Разработка веб-системы онлайн-бронирования\n"
                "групповых транзитов «ТываТранзит»"
            )
            p.font.size = Pt(28)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER

    info = slide.shapes.add_textbox(Inches(1.2), Inches(4.6), Inches(10.5), Inches(2.2))
    tf = info.text_frame
    tf.word_wrap = True
    blocks = [
        ("Специальность: 09.02.07 «Информационные системы и программирование»", False),
        (f"Студент: {STUDENT}, {GROUP}", False),
        (SUPERVISOR, False),
        ("Кызыл, 2026", True),
    ]
    for i, (text, bold) in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(18)
        p.font.bold = bold
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)


def main() -> None:
    prs = Presentation(str(TEMPLATE))

    # Slide 1 — title
    fill_title_slide(prs.slides[0])

    # Slide 2 — goal & tasks
    slide = prs.slides[1]
    set_title(slide, "Цель и задачи")
    add_bullets(
        slide,
        [
            "Цель: разработать веб-систему «ТываТранзит» для онлайн-бронирования "
            "транзитов из Кызыла на озёра Тывы",
            "Проанализировать деятельность заказчика и процесс бронирования",
            "Сформулировать требования к системе",
            "Спроектировать архитектуру клиент–серверного приложения",
            "Реализовать серверную часть (Django REST) и клиентский интерфейс",
            "Развернуть систему на Railway и провести тестирование",
        ],
    )

    # Slide 3 — client
    slide = prs.slides[2]
    set_title(slide, "О предприятии — заказчике")
    add_bullets(
        slide,
        [
            "ИП Ширин-оол Долаана Май-ооловна, г. Кызыл, Республика Тыва",
            "Вид деятельности: организация пассажирских транзитов (не туроператор)",
            "Маршруты: Кызыл → Чагытай, Тере-Холь, Азас; обратные выезды",
            "Пункт сбора: ул. Кочетова, 2",
            "Парк: минивэн (до 8 мест) и микроавтобус (до 12 мест)",
            "До автоматизации: бронирование по телефону и в мессенджерах",
        ],
    )

    # Slide 4 — relevance
    slide = prs.slides[3]
    set_title(slide, "Актуальность")
    add_bullets(
        slide,
        [
            "Озёра Тывы популярны у туристов, спрос на транзиты растёт в сезон",
            "Ручной учёт мест приводит к ошибкам и перебронированию",
            "Нет единого публичного расписания — клиенты звонят повторно",
            "Заявки на личные поездки теряются в переписке",
            "Небольшому перевозчику не нужна сложная CRM — достаточно "
            "специализированного веб-сервиса",
        ],
    )

    # Slide 5 — technologies
    slide = prs.slides[4]
    set_title(slide, "Используемые технологии")
    add_bullets(
        slide,
        [
            "Backend: Python 3.12, Django 5, Django REST Framework",
            "База данных: PostgreSQL (прод), SQLite (локальная разработка)",
            "Frontend: HTML5, CSS3, JavaScript (многостраничный интерфейс)",
            "Статика: WhiteNoise; контейнеризация: Docker",
            "Деплой: Railway, автосборка из GitHub (ветка main)",
            "Уведомления: SMTP (email-билеты пассажиру и перевозчику)",
            "Репозиторий: github.com/bailak-bz/TyvaTransit",
        ],
    )

    # Slide 6 — design part 1
    slide = prs.slides[5]
    set_title(slide, "Дизайн, функционал приложения")
    add_body_text(
        slide,
        [
            "Публичная часть сайта:",
            "• Главная — выбор типа поездки (общий рейс / личная)",
            "• Каталог рейсов с фильтрацией по направлению и дате",
            "• Бронирование мест, проверка билета по номеру и телефону",
        ],
        left=Inches(0.5),
        top=Inches(1.8),
        width=Inches(5.2),
        height=Inches(4.5),
    )
    img = SCREENSHOTS / "index.png"
    if img.exists():
        add_image(slide, img, left=Inches(6.0), top=Inches(1.9), width=Inches(6.3))

    # Slide 7 — design part 2
    slide = prs.slides[6]
    set_title(slide, "Дизайн, функционал приложения")
    add_body_text(
        slide,
        [
            "API и администрирование:",
            "• REST API: /api/trips/, /api/bookings/shared/, /api/bookings/private/",
            "• Уникальный код брони: TYV-T-ГГГГ-XXXX (общий), TYV-P-... (личный)",
            "• Django Admin — управление направлениями, рейсами и бронями",
            "• Транзакции и блокировка мест при одновременной покупке",
        ],
        left=Inches(0.5),
        top=Inches(1.8),
        width=Inches(5.2),
        height=Inches(4.5),
    )
    img = SCREENSHOTS / "trips.png"
    if img.exists():
        add_image(slide, img, left=Inches(6.0), top=Inches(1.9), width=Inches(6.3))

    # Slide 8 — conclusion
    slide = prs.slides[7]
    set_title(slide, "Заключение")
    add_bullets(
        slide,
        [
            "Создана работающая веб-система бронирования транзитов «ТываТранзит»",
            "Автоматизированы общие рейсы и заявки на личные поездки",
            "Пассажир получает билет на email и может проверить бронь онлайн",
            "Система развёрнута на Railway (Django + PostgreSQL)",
            "Перспективы: реальная оплата (СБП), SMS, личный кабинет перевозчика",
        ],
    )

    # Slide 9 — thanks
    slide = prs.slides[8]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text == "Ссылка на ресурс для демонстрации":
            shape.text_frame.text = DEMO_URL
            for p in shape.text_frame.paragraphs:
                p.font.size = Pt(22)
                p.alignment = PP_ALIGN.CENTER

    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    main()
